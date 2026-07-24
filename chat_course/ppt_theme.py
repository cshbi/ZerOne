# -*- coding: utf-8 -*-
"""
공용 PPT 테마 (Deck)
------------------------------------------------------------
주차별 make_ppt.py 가 공유하는 슬라이드 생성기.
사용법:
    from ppt_theme import Deck
    d = Deck()
    d.title(kicker=..., title=..., subtitle=..., notes=...)
    d.big_question(question=..., hint=..., notes=...)
    d.bullets("제목", ["항목", ("들여쓴 항목", 1), ...], notes=...)
    d.code("제목", "코드 문자열", caption=..., notes=...)
    d.terminals("제목", [("패널 제목", "패널 내용"), ...], notes=...)
    d.takeaway(headline=..., points=[...], notes=...)
    d.save("파일.pptx")

디자인 규칙(모든 주차 공통):
  - 16:9 (13.333 x 7.5 in), 본문 맑은 고딕 / 코드·터미널 Consolas
  - 표지·질문·정리 슬라이드에만 왼쪽 파란 세로 바
  - 내용 슬라이드는 제목(네이비 30pt bold) + 파란 밑줄 바
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# ── 팔레트 ──────────────────────────────────────────────
NAVY = RGBColor(0x1F, 0x3A, 0x5F)      # 진한 배경, 제목 글자
ACCENT = RGBColor(0x2E, 0x86, 0xC1)    # 포인트 파랑 (바, 킥커, 패널 제목)
LIGHT_BG = RGBColor(0xF1, 0xF5, 0xF9)  # 질문 슬라이드 배경
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BODY = RGBColor(0x22, 0x29, 0x33)      # 본문 불릿 글자
MUTED = RGBColor(0x6B, 0x7A, 0x8C)     # 보조 설명, 캡션
SOFT = RGBColor(0xC7, 0xD3, 0xE0)      # 어두운 배경 위 보조 글자
CODE_BG = RGBColor(0x1E, 0x1E, 0x2E)   # 코드/터미널 상자
CODE_FG = RGBColor(0xEA, 0xEC, 0xF0)
CODE_COMMENT = RGBColor(0x7C, 0x8A, 0x9E)

FONT = "맑은 고딕"
MONO = "Consolas"


class Deck:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        self._blank = self.prs.slide_layouts[6]

    # ── 내부 도우미 ─────────────────────────────────────
    def _slide(self, bg):
        s = self.prs.slides.add_slide(self._blank)
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = bg
        return s

    def _notes(self, slide, notes):
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    def _rect(self, slide, x, y, w, h, color):
        from pptx.enum.shapes import MSO_SHAPE
        r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(x), Inches(y), Inches(w), Inches(h))
        r.fill.solid()
        r.fill.fore_color.rgb = color
        r.line.fill.background()
        return r

    def _sidebar(self, slide):
        """표지·질문·정리 슬라이드의 왼쪽 파란 세로 바."""
        self._rect(slide, 0, 0, 0.35, 7.5, ACCENT)

    def _textbox(self, slide, x, y, w, h):
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tb.text_frame.word_wrap = True
        return tb.text_frame

    def _put(self, tf, lines, size, color, bold=False, font=FONT,
             space_after=2, first=True):
        """여러 줄 텍스트를 한 상자에 넣는다. lines: str 또는 [str]."""
        if isinstance(lines, str):
            lines = lines.split("\n")
        for line in lines:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.line_spacing = 1.0
            p.space_after = Pt(space_after)
            r = p.add_run()
            r.text = line if line else " "
            r.font.name = font
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color
        return tf

    def _heading(self, slide, title):
        """내용 슬라이드 공통: 제목 + 파란 밑줄 바."""
        tf = self._textbox(slide, 0.70, 0.45, 12.00, 0.90)
        self._put(tf, title, 30, NAVY, bold=True)
        self._rect(slide, 0.75, 1.35, 2.20, 0.07, ACCENT)

    # ── 슬라이드 종류 ───────────────────────────────────
    def title(self, kicker, title, subtitle, notes=""):
        s = self._slide(NAVY)
        self._sidebar(s)
        self._put(self._textbox(s, 1.00, 2.10, 11.50, 0.60), kicker, 18, ACCENT, bold=True)
        self._put(self._textbox(s, 1.00, 2.70, 11.50, 2.00), title, 44, WHITE, bold=True)
        self._put(self._textbox(s, 1.00, 4.70, 11.50, 1.00), subtitle, 20, SOFT)
        self._notes(s, notes)

    def big_question(self, question, hint="", notes=""):
        s = self._slide(LIGHT_BG)
        self._sidebar(s)
        self._put(self._textbox(s, 1.10, 1.00, 11.00, 0.70),
                  "이런 게 궁금하지 않나요?", 18, ACCENT, bold=True)
        self._put(self._textbox(s, 1.10, 2.40, 11.00, 2.60), question, 36, NAVY, bold=True)
        if hint:
            self._put(self._textbox(s, 1.10, 5.20, 11.00, 1.20), hint, 20, MUTED)
        self._notes(s, notes)

    def bullets(self, title, items, notes=""):
        """items: "글" (기본) 또는 ("글", 1) (들여쓴 보조 항목)."""
        s = self._slide(WHITE)
        self._heading(s, title)
        tf = self._textbox(s, 0.85, 1.70, 11.60, 5.20)
        first = True
        for item in items:
            text, level = item if isinstance(item, tuple) else (item, 0)
            if level == 0:
                self._put(tf, "•  " + text, 22, BODY, bold=True,
                          space_after=12, first=first)
            else:
                self._put(tf, "–  " + text, 19, MUTED, space_after=12, first=first)
            first = False
        self._notes(s, notes)

    def code(self, title, code, caption="", notes=""):
        s = self._slide(WHITE)
        self._heading(s, title)
        box = self._rect(s, 0.75, 1.70, 11.80, 4.90, CODE_BG)
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = Inches(0.3)
        tf.margin_top = tf.margin_bottom = Inches(0.2)
        from pptx.enum.text import MSO_ANCHOR
        tf.vertical_anchor = MSO_ANCHOR.TOP
        lines = code.split("\n")
        size = 15 if len(lines) <= 17 else 13
        first = True
        for line in lines:
            color = CODE_COMMENT if line.lstrip().startswith("#") else CODE_FG
            self._put(tf, line, size, color, font=MONO, first=first)
            first = False
        if caption:
            self._put(self._textbox(s, 0.80, 6.75, 11.80, 0.55), caption, 16, MUTED)
        self._notes(s, notes)

    def terminals(self, title, panels, notes=""):
        """panels: [(패널 제목, 패널 내용), ...] — 실제 실행 결과를 캡처해 넣는다."""
        s = self._slide(WHITE)
        self._heading(s, title)
        n = len(panels)
        margin, gap = 0.55, 0.30
        w = (13.333 - 2 * margin - (n - 1) * gap) / n
        for i, (head, body) in enumerate(panels):
            x = margin + i * (w + gap)
            hbox = self._rect(s, x, 1.70, w, 0.50, ACCENT)
            htf = hbox.text_frame
            htf.margin_left = htf.margin_right = Inches(0.15)
            self._put(htf, head, 13, WHITE, bold=True, font=MONO)
            bbox = self._rect(s, x, 2.25, w, 4.70, CODE_BG)
            btf = bbox.text_frame
            btf.word_wrap = True
            btf.margin_left = btf.margin_right = Inches(0.3)
            btf.margin_top = btf.margin_bottom = Inches(0.2)
            from pptx.enum.text import MSO_ANCHOR
            btf.vertical_anchor = MSO_ANCHOR.TOP
            self._put(btf, body, 13, CODE_FG, font=MONO)
        self._notes(s, notes)

    def takeaway(self, headline, points, notes=""):
        s = self._slide(NAVY)
        self._sidebar(s)
        self._put(self._textbox(s, 1.00, 0.90, 11.30, 0.70),
                  "오늘 손에 남는 것", 18, ACCENT, bold=True)
        self._put(self._textbox(s, 1.00, 2.00, 11.30, 2.20), headline, 34, WHITE, bold=True)
        tf = self._textbox(s, 1.00, 4.40, 11.30, 2.60)
        first = True
        for pt_ in points:
            self._put(tf, "•  " + pt_, 20, SOFT, space_after=8, first=first)
            first = False
        self._notes(s, notes)

    def save(self, path):
        self.prs.save(path)
